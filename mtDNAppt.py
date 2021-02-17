#author:Zelin Li
#date:2020/02/29
#utility:generating ppt from genomic analysis results.
#documentation:https://python-pptx.readthedocs.io/en/latest/
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]

slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Pg mtDNA assembly"
subtitle.text = "Zelin Li 2020/04/10"

def text_read(filename):
    try:
        file = open(filename,'r')
    except IOError:
        error = []
        return error
    strand = file.readlines()
    for i in range(len(strand)):
        strand[i] = strand[i][:len(strand[i])-1]
    file.close()
    return strand

test_strand = text_read('list.txt')

for j in range(len(test_strand)):
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	title.text = test_strand[j]

	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	title_shape = shapes.title
	body_shape = shapes.placeholders[1]
	title_shape.text = "abyss"
	tf = body_shape.text_frame
	with open(test_strand[j]+"_a_scaffolds.fa_mt_spp.sel.info.txt",'r') as file1:
		mtlength = file1.read()
	tf.text = mtlength
	with open(test_strand[j]+"_a_stats.txt",'r') as file2:
		stats = file2.read()
	p = tf.add_paragraph()
	p.text = stats

	slide = prs.slides.add_slide(blank_slide_layout)
	left = Inches(1.3)
	top = Inches(0)
	pic = slide.shapes.add_picture(test_strand[j]+"_a_scaffolds.fa_mt.png", left, top, height = Inches(7.5))

	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	title_shape = shapes.title
	body_shape = shapes.placeholders[1]
	title_shape.text = "platanus"
	tf = body_shape.text_frame
	with open(test_strand[j]+"_p_scaffolds.fa_mt_spp.sel.info.txt",'r') as file3:
		mtlength = file3.read()
	tf.text = mtlength
	with open(test_strand[j]+"_p_stats.txt",'r') as file4:
		stats = file4.read()
	p = tf.add_paragraph()
	p.text = stats

	slide = prs.slides.add_slide(blank_slide_layout)
	left = Inches(1.3)
	top = Inches(0)
	pic = slide.shapes.add_picture(test_strand[j]+"_p_scaffolds.fa_mt.png", left, top, height = Inches(7.5))

	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	title_shape = shapes.title
	body_shape = shapes.placeholders[1]
	title_shape.text = "spades"
	tf = body_shape.text_frame
	with open(test_strand[j]+"_s_scaffolds.fa_mt_spp.sel.info.txt",'r') as file5:
		mtlength = file5.read()
	tf.text = mtlength
	with open(test_strand[j]+"_s_stats.txt",'r') as file6:
		stats = file6.read()
	p = tf.add_paragraph()
	p.text = stats

	slide = prs.slides.add_slide(blank_slide_layout)
	left = Inches(1.3)
	top = Inches(0)
	pic = slide.shapes.add_picture(test_strand[j]+"_s_scaffolds.fa_mt.png", left, top, height = Inches(7.5))

pass

prs.save('Pg_mtDNA.pptx')
